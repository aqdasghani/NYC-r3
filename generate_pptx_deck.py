import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    COLOR_BG_DARK = RGBColor(6, 49, 32)      # Deep Emerald (#063120)
    COLOR_PRIMARY = RGBColor(15, 169, 88)    # Vibrant Green (#0FA958)
    COLOR_TEXT_DARK = RGBColor(15, 23, 42)   # Slate (#0F172A)
    COLOR_TEXT_LIGHT = RGBColor(248, 250, 252) # White Slate
    COLOR_CARD_BG = RGBColor(240, 253, 244)  # Light Green Tint (#F0FDF4)
    COLOR_CARD_BORDER = RGBColor(187, 247, 208)

    blank_layout = prs.slide_layouts[6]

    # Helper: Add Slide Header
    def add_slide_header(slide, title_text, category_text="GREEN QUANT AI — HACKATHON DECK"):
        # Header banner
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_PRIMARY
        p_cat.font.name = 'Inter'
        
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_BG_DARK
        p_title.font.name = 'Inter'

    # SLIDE 1: Title Slide
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_BG_DARK
    bg1.line.fill.background()

    tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "🌿 GREEN QUANT AI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.font.name = 'Inter'
    
    p2 = tf1.add_paragraph()
    p2.text = "Smart Retail. Zero Waste. Stop Money Walking Out."
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLOR_TEXT_LIGHT
    p2.font.name = 'Inter'
    p2.space_before = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = "Team NYC-r3  |  Hackathon Judgment Presentation  |  https://greenshop-ai.vercel.app"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(148, 163, 184)
    p3.font.name = 'Inter'
    p3.space_before = Pt(30)

    slide1.notes_slide.notes_text_frame.text = "Good morning judges! We are Team NYC-r3, and today we’re introducing Green Quant AI—a platform built to stop retail stores from losing thousands of dollars every month to expired inventory."

    # SLIDE 2: The Problem
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide2, "The $1.2 Trillion Bleed in Retail Inventory")

    # 3 Metric Cards
    metrics = [
        ("$1.2 Trillion", "Global Annual Loss", "Lost worldwide to retail inventory waste and unsold expired stock."),
        ("₹1.5 Lakh Crore", "Indian Kirana Loss", "Annual food & medicine waste across local independent store networks."),
        ("15–20% Margin", "Direct Profit Bleed", "Bottom-line net margin lost by retail owners due to manual tracking errors.")
    ]
    for i, (val, lbl, desc) in enumerate(metrics):
        x = Inches(0.8 + i * 3.9)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(3.7), Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.4)

        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.font.name = 'Inter'

        p_lbl = tf.add_paragraph()
        p_lbl.text = lbl
        p_lbl.font.size = Pt(14)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = COLOR_BG_DARK
        p_lbl.font.name = 'Inter'
        p_lbl.space_before = Pt(8)

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = COLOR_TEXT_DARK
        p_desc.font.name = 'Inter'
        p_desc.space_before = Pt(12)

    slide2.notes_slide.notes_text_frame.text = "Walk into any local kirana store or pharmacy. Items sit on shelves expiring silently. Traditional POS tools are static ledgers—they tell you what you sold yesterday, but never alert you about what is expiring tomorrow."

    # SLIDE 3: The Solution
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide3, "The Solution: Quantitative Retail Intelligence")

    pillars = [
        ("1. Smart Digitization", "Instant Camera Barcode & Invoice OCR parsing. Auto-creates stock batches with expiry dates upon receipt."),
        ("2. Automated FEFO Engine", "First-Expired, First-Out algorithm. Automatically allocates stock from the earliest expiring batch during checkout."),
        ("3. Proactive AI Actions", "Generates 4 automated decision cards (Discount, Transfer, Reorder, Sell-First) to liquidate stock before loss.")
    ]
    for i, (title, text) in enumerate(pillars):
        y = Inches(1.6 + i * 1.7)
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(1.4))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = COLOR_PRIMARY

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_BG_DARK
        p.font.name = 'Inter'

        p_body = tf.add_paragraph()
        p_body.text = text
        p_body.font.size = Pt(12)
        p_body.font.color.rgb = COLOR_TEXT_DARK
        p_body.font.name = 'Inter'
        p_body.space_before = Pt(4)

    slide3.notes_slide.notes_text_frame.text = "Green Quant AI treats every product as a financial asset with a clock ticking down. We digitize stock instantly via invoice OCR, enforce FEFO during checkout, and give owners proactive AI decision cards."

    # SLIDE 4: Core Innovation (FEFO & Cards)
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide4, "Core Innovation: FEFO Engine & AI Decision Cards")

    cards_data = [
        ("🚨 SELL FIRST", "Batch #B102 expires in 4 days. Position at checkout counter immediately."),
        ("📉 DISCOUNT", "Apply 25% discount on Batch #B104 to recover ₹4,500 capital."),
        ("🔄 TRANSFER", "Move 40 units to Store #2 where sales velocity is 3x higher."),
        ("🛒 REORDER", "Stock reaches safety threshold in 2 days. Reorder to prevent stockouts.")
    ]
    for i, (card_title, card_desc) in enumerate(cards_data):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.6 + row * 2.6)

        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = card_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.font.name = 'Inter'

        p_desc = tf.add_paragraph()
        p_desc.text = card_desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = COLOR_TEXT_DARK
        p_desc.font.name = 'Inter'
        p_desc.space_before = Pt(8)

    slide4.notes_slide.notes_text_frame.text = "Instead of giving store managers messy Excel sheets, our engine generates 4 simple action cards: Sell First, Discount, Transfer, or Reorder. One click executes the financial intervention."

    # SLIDE 5: Product Features & Smart Receiving
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide5, "Built for Ground-Level Store Operations")

    feats = [
        ("📷 Camera Invoice OCR", "Snap a photo of vendor paper invoice → AI extracts line items, batch numbers, cost price, and expiry dates automatically."),
        ("⚡ Sub-50ms Barcode Scanner", "Uses built-in camera for rapid stock entry and instant POS checkout verification."),
        ("🔄 Multi-Branch Transfers", "3-step wizard for transferring stock between main store and godowns/branches."),
        ("📑 1-Click Z-Report Export", "Instant CSV daily sales & tax audit export for accounting reconciliation.")
    ]
    for i, (title, desc) in enumerate(feats):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.6 + row * 2.6)

        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(203, 213, 225)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_BG_DARK
        p.font.name = 'Inter'

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11.5)
        p_desc.font.color.rgb = COLOR_TEXT_DARK
        p_desc.font.name = 'Inter'
        p_desc.space_before = Pt(6)

    slide5.notes_slide.notes_text_frame.text = "Retail staff don't have time to type 500 items manually. With our platform, staff take a photo of a vendor paper invoice. Our OCR engine parses line items, batch numbers, and expiry dates automatically."

    # SLIDE 6: Architecture
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide6, "Production-Grade Hybrid Offline Architecture")

    arch_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0))
    arch_box.fill.solid()
    arch_box.fill.fore_color.rgb = COLOR_CARD_BG
    arch_box.line.color.rgb = COLOR_PRIMARY

    tf = arch_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = Inches(0.4)

    p = tf.paragraphs[0]
    p.text = "NEXT.JS 16 CLIENT APP  ➔  FACADE LAYER (lib/api.ts)  ➔  FASTAPI REST / INDEXEDDB FALLBACK"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_BG_DARK

    bullets = [
        "Offline-First Resilience: Retail internet in Tier-2/3 cities is unreliable. Our Facade pattern transparently degrades to local IndexedDB storage so POS checkout never freezes.",
        "Async Python Engine: FastAPI + APScheduler executes background sweeps every 15 minutes to calculate expiry risk schedules and update Green Scores.",
        "Real-Time WebSocket Sync: Broadcasts stock updates instantly across multiple cashier terminals upon item checkout.",
        "Database Portability: Native cross-compatibility with SQLite for instant local boot and PostgreSQL for enterprise scale."
    ]
    for b in bullets:
        pb = tf.add_paragraph()
        pb.text = "• " + b
        pb.font.size = Pt(12)
        pb.font.color.rgb = COLOR_TEXT_DARK
        pb.space_before = Pt(12)

    slide6.notes_slide.notes_text_frame.text = "If the internet goes down in a kirana store, checkout cannot stop. Our Facade Architecture seamlessly switches to local IndexedDB storage, ensuring 100% operational uptime."

    # SLIDE 7: Dual AI
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide7, "Dual-Layer AI & Conversational WhatsApp Bot")

    for i, (layer, tech, desc) in enumerate([
        ("Layer 1: Deterministic AI Engines", "Python (expiry_engine.py, forecast_engine.py)", "Computes exact mathematical risk schedules, lead-time velocity, and discount calculations with ZERO hallucination."),
        ("Layer 2: Large Language Models", "Google Gemini 1.5 Pro / GPT-4o-mini", "Powers natural language WhatsApp store assistant and daily executive AI briefing summaries.")
    ]):
        x = Inches(0.8 + i * 5.9)
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(5.6), Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = COLOR_PRIMARY

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = layer
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_BG_DARK

        p_t = tf.add_paragraph()
        p_t.text = tech
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_PRIMARY
        p_t.space_before = Pt(4)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_TEXT_DARK
        p_d.space_before = Pt(12)

    slide7.notes_slide.notes_text_frame.text = "We combine deterministic mathematical engines for exact financial calculations with Gemini and GPT for conversational WhatsApp queries. Owners can literally text their store: 'What is expiring this week?'"

    # SLIDE 8: Competitor Matrix
    slide8 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide8, "Competitor Landscape & Market Advantage")

    table_shape = slide8.shapes.add_table(5, 4, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0))
    table = table_shape.table

    headers = ["Feature", "Tally / Vyapar", "Zoho Inventory", "Green Quant AI"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_BG_DARK
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_LIGHT
        p.font.size = Pt(13)

    data = [
        ["Dynamic FEFO Engine", "❌ Manual / None", "❌ Basic Date", "✅ Dynamic FEFO Allocation"],
        ["AI Risk Action Cards", "❌ None", "❌ None", "✅ Proactive Markdown/Transfer"],
        ["Invoice Photo OCR", "❌ Manual Entry", "❌ Manual Entry", "✅ Instant Camera OCR Scan"],
        ["Offline Hybrid Uptime", "Desktop Only", "❌ Fails without Cloud", "✅ 100% Offline Hybrid Uptime"]
    ]
    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11.5)
            if col_idx == 3:
                p.font.bold = True
                p.font.color.rgb = COLOR_PRIMARY

    slide8.notes_slide.notes_text_frame.text = "Legacy tools like Tally or Vyapar are built for accountants, not store operations. Green Quant AI is the only platform providing FEFO allocation, invoice OCR, and predictive AI action cards."

    # SLIDE 9: Business Model
    slide9 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide9, "Unit Economics & Monetization Roadmap")

    tiers = [
        ("Starter Tier", "FREE", "Up to 200 SKUs, basic batch tracking, manual sales."),
        ("Pro Store Plan", "₹1,499 / mo ($19)", "Unlimited SKUs, Camera Barcode, AI Action Cards, WhatsApp Assistant. 10x ROI for stores."),
        ("Enterprise Multi-Store", "₹3,999 / mo / store", "Multi-branch stock transfers, supplier scoring, dedicated DB sync, custom RBAC.")
    ]
    for i, (title, price, desc) in enumerate(tiers):
        x = Inches(0.8 + i * 3.9)
        card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(3.7), Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_BG_DARK

        p_pr = tf.add_paragraph()
        p_pr.text = price
        p_pr.font.size = Pt(22)
        p_pr.font.bold = True
        p_pr.font.color.rgb = COLOR_PRIMARY
        p_pr.space_before = Pt(8)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11.5)
        p_d.font.color.rgb = COLOR_TEXT_DARK
        p_d.space_before = Pt(12)

    slide9.notes_slide.notes_text_frame.text = "Our pricing is simple: ₹1,499 per month. For a kirana store losing ₹15,000 every month to spoiled food, Green Quant AI pays for itself within the first week."

    # SLIDE 10: Roadmap & Call to Action
    slide10 = prs.slides.add_slide(blank_layout)
    bg10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg10.fill.solid()
    bg10.fill.fore_color.rgb = COLOR_BG_DARK
    bg10.line.fill.background()

    tb10 = slide10.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(5.0))
    tf10 = tb10.text_frame
    tf10.word_wrap = True

    p = tf10.paragraphs[0]
    p.text = "Zero Waste. Infinite Scale."
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    p2 = tf10.add_paragraph()
    p2.text = "Building the sustainable retail infrastructure of tomorrow."
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_TEXT_LIGHT
    p2.space_before = Pt(8)

    p3 = tf10.add_paragraph()
    p3.text = "• Live Application: https://greenshop-ai.vercel.app\n• GitHub Repository: https://github.com/aqdasghani/NYC-r3\n• Team: NYC-r3\n\nThank you! We welcome your questions."
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.font.name = 'Inter'
    p3.space_before = Pt(24)

    slide10.notes_slide.notes_text_frame.text = "Green Quant AI isn't just about protecting retail profits—it's about stopping food and medicine waste at scale. Check out our live app at greenshop-ai.vercel.app. Thank you, and we're ready for your questions!"

    desktop_path = r"C:\Users\sbhrn\Desktop\GreenQuant_AI_Hackathon_Pitch.pptx"
    artifact_path = r"C:\Users\sbhrn\.gemini\antigravity\brain\e98b1cc8-3ff6-4015-bcb8-45ce58f23e66\GreenQuant_AI_Hackathon_Pitch.pptx"

    prs.save(desktop_path)
    prs.save(artifact_path)
    print("PowerPoint presentation generated successfully!")

if __name__ == "__main__":
    create_presentation()
